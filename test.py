def func(temp):
    from pptx import Presentation
    prs = Presentation('temp2.pptx')
    # To get shapes in your slides
    slides = [slide for slide in prs.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    replace_text(temp,shapes,prs)

def replace_text(replacements, shapes,prs):
    # print(len(shapes))
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    # print('hai')
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
    print(replacements)
    for i in replacements.values():
        prs.save("{}.pptx".format(i))
        break
 

with open("datas.txt") as file:   
    data = file.read()
    for i in data.split("\n"):
        temp=dict()
        num=0;num1=0
        for j in range(len(i)):
            if(i[j]==',' and num==0):
                temp['name']=i[:j].strip().upper()
                num=j
            elif(i[j]==',' and num!=0):
                temp['college_nme']=i[num+1:j].strip().upper()
                temp['title']=i[j+1:].strip().upper()
        # print(temp)
            # if(i[j].isnumeric() and num==0):
            #     temp['CHAITRA PATIL']=i[:j].strip()
            #     num=1
            # elif(i[j]=="%" and num==1):
            #     temp['83']=i[j-3:j].strip()
            #     num1=2
            #     num=j+1
            # elif(num==j and num1==2):
            #     temp['SKSVMACET, LAXMESHWAR']=i[j:].strip()
        func(temp)